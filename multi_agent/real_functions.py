"""
真实功能调用模块 - 智能体调用项目中的真实功能
Real Function Calls - Agents Call Real Project Functions
"""

import os
import json
import logging
from typing import Dict, List, Any, Optional
from datetime import datetime
import asyncio

logger = logging.getLogger(__name__)


class RealFunctionExecutor:
    """真实功能执行器 - 调用项目中已有的真实功能"""
    
    def __init__(self, project_root: str = None):
        self.project_root = project_root or os.getcwd()
        self.output_dir = os.path.join(self.project_root, "outputs")
        os.makedirs(self.output_dir, exist_ok=True)
    
    async def create_ppt(self, topic: str, slides: List[Dict[str, Any]]) -> Dict[str, Any]:
        """真实创建PPT文件"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            prs = Presentation()
            
            for slide_data in slides:
                slide_layout = prs.slide_layouts[0] if slide_data.get("type") == "title" else prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.get("title", "")
                
                if hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == 1:
                            placeholder.text = slide_data.get("subtitle", slide_data.get("content", [""])[0] if slide_data.get("content") else "")
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"ppt_{topic[:20]}_{timestamp}.pptx"
            filename = "".join(c for c in filename if c.isalnum() or c in "._-")
            filepath = os.path.join(self.output_dir, filename)
            
            prs.save(filepath)
            
            return {
                "success": True,
                "file_path": filepath,
                "filename": filename,
                "slides_count": len(slides),
                "message": f"PPT已创建: {filename}"
            }
        
        except ImportError:
            logger.warning("python-pptx not installed, using fallback")
            return await self._create_ppt_fallback(topic, slides)
        except Exception as e:
            logger.error(f"PPT creation error: {e}")
            return {"success": False, "error": str(e)}
    
    async def _create_ppt_fallback(self, topic: str, slides: List[Dict]) -> Dict[str, Any]:
        """PPT创建回退方案"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"ppt_{topic[:20]}_{timestamp}.json"
        filepath = os.path.join(self.output_dir, filename)
        
        ppt_data = {
            "topic": topic,
            "slides": slides,
            "created_at": datetime.now().isoformat(),
            "format": "json_structure"
        }
        
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(ppt_data, f, ensure_ascii=False, indent=2)
        
        return {
            "success": True,
            "file_path": filepath,
            "filename": filename,
            "slides_count": len(slides),
            "message": f"PPT结构已创建（JSON格式）: {filename}",
            "note": "安装 python-pptx 可生成真实PPT文件"
        }
    
    async def create_excel(self, data: Dict[str, Any], formulas: List[Dict[str, str]] = None) -> Dict[str, Any]:
        """真实创建Excel文件"""
        try:
            import openpyxl
            from openpyxl.styles import Font, Alignment
            
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "Sheet1"
            
            if isinstance(data, dict):
                if "headers" in data and "rows" in data:
                    for col, header in enumerate(data["headers"], 1):
                        ws.cell(row=1, column=col, value=header)
                        ws.cell(row=1, column=col).font = Font(bold=True)
                    
                    for row_idx, row_data in enumerate(data["rows"], 2):
                        for col_idx, value in enumerate(row_data, 1):
                            ws.cell(row=row_idx, column=col_idx, value=value)
            
            if formulas:
                for formula in formulas:
                    cell = formula.get("cell", "A1")
                    formula_str = formula.get("formula", "")
                    if cell and formula_str:
                        ws[cell] = formula_str
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"excel_{timestamp}.xlsx"
            filepath = os.path.join(self.output_dir, filename)
            
            wb.save(filepath)
            
            return {
                "success": True,
                "file_path": filepath,
                "filename": filename,
                "message": f"Excel文件已创建: {filename}"
            }
        
        except ImportError:
            logger.warning("openpyxl not installed, using fallback")
            return await self._create_excel_fallback(data, formulas)
        except Exception as e:
            logger.error(f"Excel creation error: {e}")
            return {"success": False, "error": str(e)}
    
    async def _create_excel_fallback(self, data: Dict, formulas: List = None) -> Dict[str, Any]:
        """Excel创建回退方案"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"excel_{timestamp}.json"
        filepath = os.path.join(self.output_dir, filename)
        
        excel_data = {
            "data": data,
            "formulas": formulas or [],
            "created_at": datetime.now().isoformat()
        }
        
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(excel_data, f, ensure_ascii=False, indent=2)
        
        return {
            "success": True,
            "file_path": filepath,
            "filename": filename,
            "message": f"Excel数据已创建（JSON格式）: {filename}",
            "note": "安装 openpyxl 可生成真实Excel文件"
        }
    
    async def create_word(self, title: str, content: str, sections: List[Dict] = None) -> Dict[str, Any]:
        """真实创建Word文档"""
        try:
            from docx import Document
            from docx.shared import Pt, Inches
            
            doc = Document()
            
            doc.add_heading(title, 0)
            
            if content:
                doc.add_paragraph(content)
            
            if sections:
                for section in sections:
                    doc.add_heading(section.get("title", ""), level=1)
                    if section.get("content"):
                        doc.add_paragraph(section["content"])
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"doc_{title[:20]}_{timestamp}.docx"
            filename = "".join(c for c in filename if c.isalnum() or c in "._-")
            filepath = os.path.join(self.output_dir, filename)
            
            doc.save(filepath)
            
            return {
                "success": True,
                "file_path": filepath,
                "filename": filename,
                "message": f"Word文档已创建: {filename}"
            }
        
        except ImportError:
            logger.warning("python-docx not installed, using fallback")
            return await self._create_word_fallback(title, content, sections)
        except Exception as e:
            logger.error(f"Word creation error: {e}")
            return {"success": False, "error": str(e)}
    
    async def _create_word_fallback(self, title: str, content: str, sections: List = None) -> Dict[str, Any]:
        """Word创建回退方案"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"doc_{title[:20]}_{timestamp}.md"
        filename = "".join(c for c in filename if c.isalnum() or c in "._-")
        filepath = os.path.join(self.output_dir, filename)
        
        md_content = f"# {title}\n\n{content}\n\n"
        if sections:
            for section in sections:
                md_content += f"## {section.get('title', '')}\n\n{section.get('content', '')}\n\n"
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(md_content)
        
        return {
            "success": True,
            "file_path": filepath,
            "filename": filename,
            "message": f"文档已创建（Markdown格式）: {filename}",
            "note": "安装 python-docx 可生成真实Word文件"
        }
    
    async def create_chart(self, chart_type: str, data: Dict[str, Any], title: str = "") -> Dict[str, Any]:
        """真实创建图表"""
        try:
            import matplotlib.pyplot as plt
            import matplotlib
            matplotlib.use('Agg')
            
            fig, ax = plt.subplots(figsize=(10, 6))
            
            labels = data.get("labels", [])
            values = data.get("values", [])
            
            if chart_type == "bar":
                ax.bar(labels, values)
            elif chart_type == "line":
                ax.plot(labels, values, marker='o')
            elif chart_type == "pie":
                ax.pie(values, labels=labels, autopct='%1.1f%%')
            else:
                ax.bar(labels, values)
            
            ax.set_title(title or "数据图表")
            ax.set_xlabel("类别")
            ax.set_ylabel("数值")
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"chart_{timestamp}.png"
            filepath = os.path.join(self.output_dir, filename)
            
            plt.savefig(filepath, dpi=150, bbox_inches='tight')
            plt.close()
            
            return {
                "success": True,
                "file_path": filepath,
                "filename": filename,
                "chart_type": chart_type,
                "message": f"图表已创建: {filename}"
            }
        
        except ImportError:
            logger.warning("matplotlib not installed, using fallback")
            return {
                "success": True,
                "chart_type": chart_type,
                "data": data,
                "message": "图表数据已准备（需要安装matplotlib生成图片）"
            }
        except Exception as e:
            logger.error(f"Chart creation error: {e}")
            return {"success": False, "error": str(e)}
    
    def get_output_files(self) -> List[Dict[str, Any]]:
        """获取已生成的文件列表"""
        files = []
        if os.path.exists(self.output_dir):
            for f in os.listdir(self.output_dir):
                filepath = os.path.join(self.output_dir, f)
                if os.path.isfile(filepath):
                    files.append({
                        "name": f,
                        "path": filepath,
                        "size": os.path.getsize(filepath),
                        "modified": datetime.fromtimestamp(os.path.getmtime(filepath)).isoformat()
                    })
        return files


real_executor = RealFunctionExecutor()
