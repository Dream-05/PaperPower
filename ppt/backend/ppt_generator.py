#!/usr/bin/env python3
"""
PPT智能生成核心模块

重新实现完整的PPT生成功能，包括：
- 需求解析
- 智能内容生成
- 图片自动生成
- 专业排版
- 完整PPTX输出
- 支持多种风格和布局
"""

import os
import sys
import json
import asyncio
import re
from dataclasses import dataclass, field
from enum import Enum
from typing import List, Dict, Optional, Any, Tuple
from pathlib import Path
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor

# 确保python-pptx库已安装
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.slide import PP_SLIDE_LAYOUT
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx library not installed. Install with: pip install python-pptx")

# 确保PIL库已安装
try:
    from PIL import Image as PILImage
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("Warning: PIL library not installed. Install with: pip install Pillow")


class PPTStyle(Enum):
    """PPT风格枚举"""
    TECH = "tech"
    BUSINESS = "business"
    MINIMAL = "minimal"
    CREATIVE = "creative"
    ACADEMIC = "academic"
    MODERN = "modern"
    CORPORATE = "corporate"
    STARTUP = "startup"


class PageType(Enum):
    """页面类型枚举"""
    COVER = "cover"
    TOC = "toc"
    CONTENT = "content"
    DATA = "data"
    COMPARISON = "comparison"
    TEAM = "team"
    TIMELINE = "timeline"
    GALLERY = "gallery"
    QUOTE = "quote"
    END = "end"


class ColorScheme:
    """颜色方案"""
    @staticmethod
    def get(style: PPTStyle) -> Dict[str, str]:
        schemes = {
            PPTStyle.TECH: {
                "primary": "#0A6EBD",
                "secondary": "#36B37E",
                "accent1": "#FF5630",
                "accent2": "#FFAB00",
                "background": "#FFFFFF",
                "text": "#172B4D"
            },
            PPTStyle.BUSINESS: {
                "primary": "#0078D4",
                "secondary": "#00B7C3",
                "accent1": "#707070",
                "accent2": "#00B7C3",
                "background": "#FFFFFF",
                "text": "#333333"
            },
            PPTStyle.MINIMAL: {
                "primary": "#333333",
                "secondary": "#666666",
                "accent1": "#999999",
                "accent2": "#E6E6E6",
                "background": "#FFFFFF",
                "text": "#333333"
            },
            PPTStyle.CREATIVE: {
                "primary": "#6C63FF",
                "secondary": "#FF6584",
                "accent1": "#FF9F1C",
                "accent2": "#4ECDC4",
                "background": "#FFFFFF",
                "text": "#2D3748"
            },
            PPTStyle.ACADEMIC: {
                "primary": "#003459",
                "secondary": "#007EA7",
                "accent1": "#00A8E8",
                "accent2": "#00171F",
                "background": "#FFFFFF",
                "text": "#333333"
            },
            PPTStyle.MODERN: {
                "primary": "#3A86FF",
                "secondary": "#8338EC",
                "accent1": "#FF006E",
                "accent2": "#FFBE0B",
                "background": "#FFFFFF",
                "text": "#1A1A2E"
            },
            PPTStyle.CORPORATE: {
                "primary": "#0047AB",
                "secondary": "#0073E6",
                "accent1": "#00B8D9",
                "accent2": "#5E35B1",
                "background": "#FFFFFF",
                "text": "#202124"
            },
            PPTStyle.STARTUP: {
                "primary": "#FF6B6B",
                "secondary": "#4ECDC4",
                "accent1": "#45B7D1",
                "accent2": "#FFA07A",
                "background": "#FFFFFF",
                "text": "#2D3436"
            }
        }
        return schemes.get(style, schemes[PPTStyle.MODERN])


class FontScheme:
    """字体方案"""
    @staticmethod
    def get(style: PPTStyle) -> Dict[str, str]:
        schemes = {
            PPTStyle.TECH: {
                "title": "微软雅黑",
                "body": "微软雅黑",
                "code": "Consolas"
            },
            PPTStyle.BUSINESS: {
                "title": "微软雅黑",
                "body": "微软雅黑",
                "code": "Consolas"
            },
            PPTStyle.MINIMAL: {
                "title": "微软雅黑",
                "body": "微软雅黑",
                "code": "Consolas"
            },
            PPTStyle.CREATIVE: {
                "title": "微软雅黑",
                "body": "微软雅黑",
                "code": "Consolas"
            },
            PPTStyle.ACADEMIC: {
                "title": "微软雅黑",
                "body": "微软雅黑",
                "code": "Consolas"
            },
            PPTStyle.MODERN: {
                "title": "微软雅黑",
                "body": "微软雅黑",
                "code": "Consolas"
            },
            PPTStyle.CORPORATE: {
                "title": "微软雅黑",
                "body": "微软雅黑",
                "code": "Consolas"
            },
            PPTStyle.STARTUP: {
                "title": "微软雅黑",
                "body": "微软雅黑",
                "code": "Consolas"
            }
        }
        return schemes.get(style, schemes[PPTStyle.MODERN])


@dataclass
class PPTConfig:
    """PPT配置"""
    title: str = ""
    subtitle: str = ""
    style: PPTStyle = PPTStyle.MODERN
    page_count: int = 12
    aspect_ratio: Tuple[int, int] = (16, 9)
    include_toc: bool = True
    include_end: bool = True
    auto_images: bool = True
    language: str = "zh-CN"


@dataclass
class PageInfo:
    """页面信息"""
    page_type: PageType
    title: str
    subtitle: str = ""
    content: List[str] = field(default_factory=list)
    bullets: List[str] = field(default_factory=list)
    image_path: Optional[str] = None
    image_url: Optional[str] = None
    layout: str = "default"
    notes: str = ""
    transition: str = "fade"


@dataclass
class PPTRequest:
    """PPT生成请求"""
    user_input: str
    keywords: List[str] = field(default_factory=list)
    style: PPTStyle = PPTStyle.MODERN
    page_count: Optional[int] = None
    custom_structure: Optional[List[str]] = None
    uploaded_images: List[str] = field(default_factory=list)
    user_id: str = "default"
    config: Optional[PPTConfig] = None


@dataclass
class PPTResult:
    """PPT生成结果"""
    success: bool
    file_path: Optional[str] = None
    pages: List[PageInfo] = field(default_factory=list)
    assets_used: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    error: Optional[str] = None


class IntentParser:
    """意图解析器"""
    
    STYLE_KEYWORDS = {
        PPTStyle.TECH: ["科技", "技术", "AI", "智能", "数字化", "互联网", "算法", "编程"],
        PPTStyle.BUSINESS: ["商务", "企业", "商业", "汇报", "工作", "管理", "战略"],
        PPTStyle.MINIMAL: ["极简", "简约", "简洁", "现代", "干净"],
        PPTStyle.CREATIVE: ["创意", "设计", "艺术", "时尚", "创意"],
        PPTStyle.ACADEMIC: ["学术", "论文", "研究", "教育", "教学"],
        PPTStyle.MODERN: ["现代", "时尚", "潮流", "当代"],
        PPTStyle.CORPORATE: ["企业", "公司", "集团", "组织"],
        PPTStyle.STARTUP: ["创业", "初创", "创新", "startup"]
    }
    
    STRUCTURES = {
        "项目介绍": ["封面", "目录", "项目背景", "解决方案", "技术架构", "商业模式", "团队介绍", "财务预测", "风险分析", "总结展望", "封底"],
        "工作汇报": ["封面", "目录", "工作概述", "完成情况", "成果展示", "问题分析", "解决方案", "下步计划", "总结", "封底"],
        "产品发布": ["封面", "目录", "市场痛点", "产品介绍", "核心功能", "技术优势", "用户场景", "市场策略", "发布计划", "封底"],
        "商业计划": ["封面", "目录", "执行摘要", "市场分析", "产品服务", "运营策略", "团队介绍", "财务规划", "风险评估", "融资需求", "封底"],
        "培训课件": ["封面", "目录", "学习目标", "知识讲解", "案例分析", "实践练习", "总结回顾", "Q&A", "封底"],
        "学术报告": ["封面", "目录", "研究背景", "文献综述", "研究方法", "实验结果", "数据分析", "讨论", "结论", "参考文献", "封底"]
    }
    
    def parse(self, user_input: str) -> PPTRequest:
        """解析用户输入"""
        request = PPTRequest(user_input=user_input)
        
        # 提取关键词
        request.keywords = self._extract_keywords(user_input)
        
        # 检测风格
        request.style = self._detect_style(user_input)
        
        # 预估页数
        estimated_pages = self._estimate_pages(user_input)
        if estimated_pages:
            request.page_count = estimated_pages
        
        # 生成配置
        request.config = PPTConfig(
            title=user_input[:100],
            style=request.style,
            page_count=request.page_count or 12
        )
        
        return request
    
    def _extract_keywords(self, text: str) -> List[str]:
        """提取关键词"""
        keywords = []
        
        # 提取引号、括号中的内容
        patterns = [
            r'"([^"]+)"',
            r"'([^']+)'",
            r'【([^】]+)】',
            r'《([^》]+)》',
            r'\(([^\)]+)\)'
        ]
        
        for pattern in patterns:
            matches = re.findall(pattern, text)
            keywords.extend(matches)
        
        # 提取中文关键词
        stop_words = {
            "的", "和", "与", "或", "等", "及", "一个", "这个", "那个", 
            "是", "在", "了", "有", "我", "你", "他", "她", "它",
            "要", "会", "能", "到", "以", "对", "就", "但", "而"
        }
        
        words = re.findall(r'[\u4e00-\u9fa5]{2,}|[a-zA-Z]{3,}', text)
        keywords.extend([w for w in words if w not in stop_words and len(w) >= 2])
        
        # 去重并限制数量
        unique_keywords = list(set(keywords))
        return unique_keywords[:15]  # 最多15个关键词
    
    def _detect_style(self, text: str) -> PPTStyle:
        """检测风格"""
        text_lower = text.lower()
        
        scores = {style: 0 for style in PPTStyle}
        
        for style, keywords in self.STYLE_KEYWORDS.items():
            for keyword in keywords:
                if keyword.lower() in text_lower:
                    scores[style] += 1
        
        max_score = max(scores.values())
        if max_score > 0:
            return max(scores, key=scores.get)
        
        return PPTStyle.MODERN
    
    def _estimate_pages(self, text: str) -> Optional[int]:
        """预估页数"""
        # 匹配数字+页/张/PPT
        match = re.search(r'(\d+)\s*[页张PPT]', text)
        if match:
            return int(match.group(1))
        
        # 根据内容长度估计
        text_length = len(text)
        if text_length < 50:
            return 8
        elif text_length < 200:
            return 12
        elif text_length < 500:
            return 15
        else:
            return 20
    
    def suggest_structure(self, request: PPTRequest) -> List[str]:
        """建议PPT结构"""
        text = request.user_input.lower()
        
        for topic, structure in self.STRUCTURES.items():
            if topic in text:
                return structure
        
        # 默认结构
        return [
            "封面", "目录", "背景介绍", "核心内容", 
            "详细说明", "案例分析", "数据展示", 
            "总结展望", "Q&A", "封底"
        ]


class ContentGenerator:
    """内容生成器"""
    
    CONTENT_TEMPLATES = {
        "背景介绍": [
            "行业现状与发展趋势分析",
            "市场需求与痛点识别",
            "本项目的定位与价值",
            "相关技术/产品发展历程"
        ],
        "核心内容": [
            "核心概念与原理",
            "关键特性与优势",
            "技术架构与实现方案",
            "创新点与差异化"
        ],
        "详细说明": [
            "功能模块详解",
            "技术实现细节",
            "使用方法与流程",
            "性能与可靠性"
        ],
        "案例分析": [
            "成功案例展示",
            "实际应用场景",
            "用户反馈与评价",
            "经验总结与启示"
        ],
        "数据展示": [
            "关键数据指标分析",
            "趋势与对比分析",
            "预测与展望",
            "数据驱动的决策"
        ],
        "总结展望": [
            "项目成果总结",
            "未来发展规划",
            "面临的挑战与机遇",
            "合作与发展建议"
        ],
        "团队介绍": [
            "核心团队成员",
            "专业背景与经验",
            "团队优势与文化",
            "发展历程与成就"
        ],
        "商业模式": [
            "盈利模式设计",
            "市场定位与策略",
            "运营与增长计划",
            "竞争优势分析"
        ],
        "财务预测": [
            "收入预测模型",
            "成本结构分析",
            "投资回报预期",
            "风险评估与应对"
        ]
    }
    
    def __init__(self, model_manager=None):
        self.model_manager = model_manager
    
    async def generate_content(self, section: str, request: PPTRequest) -> List[str]:
        """生成内容"""
        # 尝试使用模型管理器生成智能内容
        if self.model_manager:
            try:
                prompt = f"为PPT幻灯片生成关于'{section}'的内容要点，主题是'{request.user_input}'。请提供4个具体、专业的要点，每个要点用简洁的中文句子表达。"
                response = await self.model_manager.generate_text(prompt, max_tokens=500)
                
                if response:
                    # 解析生成的内容
                    lines = response.strip().split('\n')
                    points = [line.strip().lstrip('1234567890.、）） ') for line in lines if line.strip()]
                    if points and len(points) >= 4:
                        return points[:4]
            except Exception as e:
                print(f"Content generation failed: {e}")
        
        # 回退到模板内容
        return self.CONTENT_TEMPLATES.get(section, [
            f"{section}要点一",
            f"{section}要点二",
            f"{section}要点三",
            f"{section}要点四"
        ])
    
    async def generate_image_content(self, section: str, request: PPTRequest, image_descriptions: List[str] = None) -> List[str]:
        """根据图片元素生成内容"""
        if self.model_manager and image_descriptions:
            try:
                image_context = "、".join(image_descriptions[:3])
                prompt = f"""为PPT幻灯片生成关于'{section}'的内容要点。
主题是'{request.user_input}'。
相关图片描述：{image_context}

请根据图片内容和主题，生成4个具体、专业的要点，每个要点用简洁的中文句子表达。要点应该与图片内容相关联。"""
                
                response = await self.model_manager.generate_text(prompt, max_tokens=500)
                
                if response:
                    lines = response.strip().split('\n')
                    points = [line.strip().lstrip('1234567890.、）） ') for line in lines if line.strip()]
                    if points and len(points) >= 4:
                        return points[:4]
            except Exception as e:
                print(f"Image content generation failed: {e}")
        
        return await self.generate_content(section, request)
    
    async def generate_image_caption(self, image_description: str, section: str, request: PPTRequest) -> str:
        """为图片生成说明文字"""
        if self.model_manager:
            try:
                prompt = f"""为PPT中的图片生成一句简短的说明文字。
主题：{request.user_input}
章节：{section}
图片描述：{image_description}

请生成一句15-30字的中文说明文字，简洁专业，与图片内容相关。"""
                
                response = await self.model_manager.generate_text(prompt, max_tokens=100)
                
                if response:
                    return response.strip().strip('"\'')
            except Exception as e:
                print(f"Caption generation failed: {e}")
        
        return image_description[:30] if image_description else ""
    
    def generate_title(self, section: str, request: PPTRequest) -> str:
        title = ""
        if section == "封面":
            title = request.user_input[:30]
        elif section == "目录":
            title = "目录"
        elif section == "封底":
            title = "感谢观看"
        else:
            title = section[:15] if len(section) > 15 else section
        return title
    
    def generate_subtitle(self, section: str, request: PPTRequest) -> str:
        """生成副标题"""
        subtitles = {
            "封面": "专业演示文稿",
            "目录": "Contents",
            "封底": "Thank You",
            "背景介绍": "Background Introduction",
            "核心内容": "Core Content",
            "详细说明": "Detailed Explanation",
            "案例分析": "Case Analysis",
            "数据展示": "Data Presentation",
            "总结展望": "Summary & Outlook",
            "团队介绍": "Team Introduction",
            "商业模式": "Business Model",
            "财务预测": "Financial Forecast"
        }
        return subtitles.get(section, "")
    
    async def generate_notes(self, section: str, request: PPTRequest) -> str:
        """生成讲者备注"""
        # 尝试使用模型管理器生成智能备注
        if self.model_manager:
            try:
                prompt = f"为PPT幻灯片生成关于'{section}'的讲者备注，主题是'{request.user_input}'。请提供一段专业、自然的中文备注，帮助演讲者更好地讲解内容。"
                response = await self.model_manager.generate_text(prompt, max_tokens=300)
                
                if response:
                    return response.strip()
            except Exception as e:
                print(f"Notes generation failed: {e}")
        
        # 回退到模板备注
        notes = {
            "封面": "欢迎各位，今天我将为大家介绍我们的项目。",
            "目录": "这是本次演示的主要内容。",
            "背景介绍": "这里将分析行业现状和市场需求。",
            "核心内容": "重点介绍我们的核心产品/技术。",
            "详细说明": "深入讲解技术实现和功能细节。",
            "案例分析": "通过实际案例展示产品价值。",
            "数据展示": "用数据说话，展示项目成果。",
            "总结展望": "总结项目成果并展望未来发展。",
            "Q&A": "欢迎大家提问和讨论。",
            "封底": "感谢大家的聆听和支持。"
        }
        return notes.get(section, f"【{section}】请根据实际情况补充详细讲解内容。")


class ImageGenerator:
    """图片生成器"""
    
    def __init__(self, output_dir: Optional[Path] = None, openai_api_key: Optional[str] = None):
        self.output_dir = output_dir or Path("data/assets/ppt_images")
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.openai_api_key = openai_api_key
        self.client = None
        
        # 初始化OpenAI客户端
        if self.openai_api_key:
            try:
                from openai import OpenAI
                self.client = OpenAI(api_key=self.openai_api_key)
            except ImportError:
                print("OpenAI library not installed")
            except Exception as e:
                print(f"Failed to initialize OpenAI client: {e}")
    
    async def generate_image(self, prompt: str, style: PPTStyle) -> Optional[str]:
        """生成图片"""
        try:
            import random
            image_id = f"image_{random.randint(1000, 9999)}"
            image_path = self.output_dir / f"{image_id}.png"
            
            # 尝试使用OpenAI DALL-E生成图片
            if self.client:
                try:
                    response = self.client.images.generate(
                        model="dall-e-3",
                        prompt=self._enhance_prompt(prompt, style),
                        size="1024x1024",
                        quality="standard",
                        n=1,
                    )
                    
                    # 下载图片
                    import httpx
                    async with httpx.AsyncClient() as client:
                        img_response = await client.get(response.data[0].url)
                        with open(image_path, "wb") as f:
                            f.write(img_response.content)
                    
                    return str(image_path)
                except Exception as e:
                    print(f"DALL-E generation failed: {e}")
                    # 回退到占位图像
            
            # 创建占位图像
            if PIL_AVAILABLE:
                from PIL import Image, ImageDraw, ImageFont
                img = Image.new('RGB', (1024, 1024), color=ColorScheme.get(style)["background"])
                d = ImageDraw.Draw(img)
                
                # 添加文字
                try:
                    font = ImageFont.truetype("arial.ttf", 24)
                except:
                    font = ImageFont.load_default()
                
                text = f"Generated for: {prompt[:50]}"
                text_width = d.textlength(text, font=font)
                d.text(
                    ((1024 - text_width) // 2, 500),
                    text,
                    fill=ColorScheme.get(style)["text"],
                    font=font
                )
                
                img.save(image_path)
                return str(image_path)
            else:
                # 创建空文件作为占位
                image_path.touch()
                return str(image_path)
        except Exception as e:
            print(f"Image generation failed: {e}")
            return None
    
    def _enhance_prompt(self, prompt: str, style: PPTStyle) -> str:
        """增强提示词"""
        style_descriptions = {
            PPTStyle.TECH: "modern technology, futuristic, clean design, professional",
            PPTStyle.BUSINESS: "corporate, professional, clean, business presentation",
            PPTStyle.MINIMAL: "minimalist, clean, simple, elegant design",
            PPTStyle.CREATIVE: "creative, artistic, vibrant, modern design",
            PPTStyle.ACADEMIC: "academic, educational, professional, informative",
            PPTStyle.MODERN: "modern, sleek, professional, clean design",
            PPTStyle.CORPORATE: "corporate, professional, formal, business-like",
            PPTStyle.STARTUP: "startup, innovative, modern, dynamic design"
        }
        
        style_desc = style_descriptions.get(style, "professional presentation")
        return f"{prompt}, {style_desc}, high quality, suitable for PowerPoint presentation, clear, professional"
    
    def set_api_key(self, api_key: str):
        """设置API密钥"""
        self.openai_api_key = api_key
        try:
            from openai import OpenAI
            self.client = OpenAI(api_key=api_key)
        except Exception as e:
            print(f"Failed to initialize OpenAI client: {e}")
    
    async def generate_images_for_pages(self, pages: List[PageInfo], style: PPTStyle) -> List[str]:
        """为页面生成图片"""
        generated_images = []
        
        async def generate_for_page(page):
            if page.page_type in [PageType.COVER, PageType.CONTENT, PageType.GALLERY]:
                prompt = f"{page.title} {page.subtitle}"
                image_path = await self.generate_image(prompt, style)
                if image_path:
                    page.image_path = image_path
                    generated_images.append(image_path)
        
        tasks = [generate_for_page(page) for page in pages]
        await asyncio.gather(*tasks)
        
        return generated_images


class LayoutEngine:
    """布局引擎"""
    
    LAYOUTS = {
        PageType.COVER: "cover",
        PageType.TOC: "toc",
        PageType.CONTENT: "content",
        PageType.DATA: "data",
        PageType.COMPARISON: "comparison",
        PageType.TEAM: "team",
        PageType.TIMELINE: "timeline",
        PageType.GALLERY: "gallery",
        PageType.QUOTE: "quote",
        PageType.END: "end"
    }
    
    def get_layout(self, page_type: PageType) -> str:
        """获取布局"""
        return self.LAYOUTS.get(page_type, "content")
    
    def get_smart_layout(self, page_type: PageType, has_image: bool = False, image_count: int = 0) -> str:
        """智能选择布局"""
        base_layout = self.LAYOUTS.get(page_type, "content")
        
        if not has_image:
            return f"{base_layout}_text_only"
        
        if image_count == 1:
            return f"{base_layout}_single_image"
        elif image_count == 2:
            return f"{base_layout}_dual_image"
        elif image_count >= 3:
            return f"{base_layout}_gallery"
        
        return base_layout
    
    def calculate_position(self, layout: str, slide_width: float, slide_height: float) -> Dict[str, Tuple[float, float, float, float]]:
        """计算元素位置"""
        positions = {
            "cover": {
                "title": (Inches(1), Inches(3), Inches(11), Inches(2)),
                "subtitle": (Inches(1), Inches(5.5), Inches(11), Inches(1)),
                "image": (Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            },
            "cover_single_image": {
                "title": (Inches(1), Inches(3), Inches(11), Inches(2)),
                "subtitle": (Inches(1), Inches(5.5), Inches(11), Inches(1)),
                "image": (Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            },
            "content": {
                "title": (Inches(1), Inches(0.5), Inches(11), Inches(1)),
                "content": (Inches(1), Inches(2), Inches(6), Inches(5)),
                "image": (Inches(7.5), Inches(2), Inches(5), Inches(5))
            },
            "content_text_only": {
                "title": (Inches(1), Inches(0.5), Inches(11), Inches(1)),
                "content": (Inches(1), Inches(2), Inches(11), Inches(5))
            },
            "content_single_image": {
                "title": (Inches(1), Inches(0.5), Inches(11), Inches(1)),
                "content": (Inches(1), Inches(2), Inches(6), Inches(5)),
                "image": (Inches(7.5), Inches(2), Inches(5), Inches(5))
            },
            "content_dual_image": {
                "title": (Inches(1), Inches(0.5), Inches(11), Inches(1)),
                "content": (Inches(1), Inches(2), Inches(5), Inches(5)),
                "image1": (Inches(6.5), Inches(2), Inches(3), Inches(2.5)),
                "image2": (Inches(9.5), Inches(2), Inches(3), Inches(2.5))
            },
            "content_gallery": {
                "title": (Inches(1), Inches(0.5), Inches(11), Inches(1)),
                "content": (Inches(1), Inches(2), Inches(4), Inches(5)),
                "image1": (Inches(5.5), Inches(2), Inches(2.5), Inches(2.5)),
                "image2": (Inches(8), Inches(2), Inches(2.5), Inches(2.5)),
                "image3": (Inches(10.5), Inches(2), Inches(2.5), Inches(2.5))
            },
            "gallery": {
                "title": (Inches(1), Inches(0.5), Inches(11), Inches(1)),
                "image1": (Inches(1), Inches(2), Inches(3.5), Inches(2.5)),
                "image2": (Inches(5), Inches(2), Inches(3.5), Inches(2.5)),
                "image3": (Inches(9), Inches(2), Inches(3.5), Inches(2.5)),
                "caption": (Inches(1), Inches(5), Inches(11), Inches(1))
            },
            "toc": {
                "title": (Inches(1), Inches(0.5), Inches(11), Inches(1)),
                "content": (Inches(2), Inches(2), Inches(10), Inches(5))
            },
            "end": {
                "title": (Inches(1), Inches(3), Inches(11), Inches(2)),
                "subtitle": (Inches(1), Inches(5.5), Inches(11), Inches(1))
            }
        }
        return positions.get(layout, positions["content"])


class PPTGenerator:
    """PPT生成器"""
    
    def __init__(self, output_dir: Optional[Path] = None, model_manager=None):
        self.output_dir = output_dir or Path("output/ppt")
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        self.model_manager = model_manager
        
        self.intent_parser = IntentParser()
        self.content_generator = ContentGenerator(model_manager=model_manager)
        self.image_generator = ImageGenerator()
        self.layout_engine = LayoutEngine()
    
    async def generate(self, request: PPTRequest) -> PPTResult:
        """生成PPT"""
        try:
            # 解析请求
            if not request.keywords:
                request = self.intent_parser.parse(request.user_input)
            
            # 生成结构
            structure = self.intent_parser.suggest_structure(request)
            
            # 创建页面
            pages = await self._create_pages(request, structure)
            
            # 生成图片
            if request.config and request.config.auto_images:
                generated_images = await self.image_generator.generate_images_for_pages(pages, request.style)
            else:
                generated_images = []
            
            # 保存PPT
            file_path = await self._save_pptx(pages, request)
            
            return PPTResult(
                success=True,
                file_path=str(file_path),
                pages=pages,
                assets_used=generated_images,
                metadata={
                    "style": request.style.value,
                    "page_count": len(pages),
                    "generated_at": datetime.now().isoformat(),
                    "images_generated": len(generated_images)
                }
            )
            
        except Exception as e:
            import traceback
            traceback.print_exc()
            return PPTResult(
                success=False,
                error=str(e)
            )
    
    async def _create_pages(self, request: PPTRequest, structure: List[str]) -> List[PageInfo]:
        """创建页面"""
        pages = []
        
        # 获取上传图片的描述
        image_descriptions = []
        if request.uploaded_images:
            for img_path in request.uploaded_images:
                # 从文件名或路径中提取描述
                img_name = Path(img_path).stem
                image_descriptions.append(img_name)
        
        for i, section in enumerate(structure):
            page_type = self._determine_page_type(i, len(structure), section)
            layout = self.layout_engine.get_layout(page_type)
            
            # 根据图片生成内容（如果有图片）
            if image_descriptions and page_type in [PageType.CONTENT, PageType.GALLERY]:
                content = await self.content_generator.generate_image_content(
                    section, request, image_descriptions
                )
            else:
                content = await self.content_generator.generate_content(section, request)
            
            notes = await self.content_generator.generate_notes(section, request)
            
            # 为有图片的页面分配图片
            image_path = None
            if request.uploaded_images and page_type in [PageType.COVER, PageType.CONTENT, PageType.GALLERY]:
                img_index = min(i, len(request.uploaded_images) - 1)
                if img_index >= 0:
                    image_path = request.uploaded_images[img_index]
            
            page = PageInfo(
                page_type=page_type,
                title=self.content_generator.generate_title(section, request),
                subtitle=self.content_generator.generate_subtitle(section, request),
                content=content,
                bullets=content,
                image_path=image_path,
                layout=layout,
                notes=notes
            )
            
            pages.append(page)
        
        return pages
    
    def _determine_page_type(self, index: int, total: int, section: str) -> PageType:
        """确定页面类型"""
        if index == 0:
            return PageType.COVER
        elif index == 1:
            return PageType.TOC
        elif index == total - 1:
            return PageType.END
        elif "团队" in section:
            return PageType.TEAM
        elif "数据" in section or "图表" in section:
            return PageType.DATA
        elif "对比" in section:
            return PageType.COMPARISON
        elif "时间" in section:
            return PageType.TIMELINE
        elif "案例" in section:
            return PageType.GALLERY
        else:
            return PageType.CONTENT
    
    async def _save_pptx(self, pages: List[PageInfo], request: PPTRequest) -> Path:
        """保存PPTX文件"""
        if not PPTX_AVAILABLE:
            # 保存为JSON结构
            file_path = self.output_dir / f"{request.keywords[0] if request.keywords else 'presentation'}_structure.json"
            structure = {
                "title": request.user_input,
                "style": request.style.value,
                "pages": [
                    {
                        "type": p.page_type.value,
                        "title": p.title,
                        "subtitle": p.subtitle,
                        "content": p.content,
                        "bullets": p.bullets,
                        "image_path": p.image_path,
                        "layout": p.layout,
                        "notes": p.notes
                    }
                    for p in pages
                ]
            }
            with open(file_path, 'w', encoding='utf-8') as f:
                json.dump(structure, f, ensure_ascii=False, indent=2)
            return file_path
        
        # 创建PPTX文件
        prs = Presentation()
        
        # 设置幻灯片尺寸 (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # 获取颜色方案
        colors = ColorScheme.get(request.style)
        fonts = FontScheme.get(request.style)
        
        for page in pages:
            # 添加幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 智能选择布局
            has_image = page.image_path and os.path.exists(page.image_path)
            smart_layout = self.layout_engine.get_smart_layout(
                page.page_type, 
                has_image=has_image, 
                image_count=1 if has_image else 0
            )
            
            # 获取布局位置
            positions = self.layout_engine.calculate_position(smart_layout, prs.slide_width, prs.slide_height)
            
            # 添加背景图片（封面页）
            if page.page_type == PageType.COVER and has_image:
                try:
                    if "image" in positions:
                        left, top, width, height = positions["image"]
                        slide.shapes.add_picture(page.image_path, left, top, width=width, height=height)
                except Exception as e:
                    print(f"Error adding background image: {e}")
            
            # 添加内容图片
            elif has_image and page.page_type != PageType.COVER:
                try:
                    if "image" in positions:
                        left, top, width, height = positions["image"]
                        slide.shapes.add_picture(page.image_path, left, top, width=width, height=height)
                except Exception as e:
                    print(f"Error adding image: {e}")
            
            # 添加标题
            if "title" in positions:
                left, top, width, height = positions["title"]
                title_box = slide.shapes.add_textbox(left, top, width, height)
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = page.title
                title_para.font.size = Pt(36)
                title_para.font.bold = True
                title_para.font.name = fonts["title"]
                title_para.font.color.rgb = self._hex_to_rgb(colors["text"])
                title_para.alignment = PP_ALIGN.CENTER
            
            # 添加副标题
            if page.subtitle and "subtitle" in positions:
                left, top, width, height = positions["subtitle"]
                subtitle_box = slide.shapes.add_textbox(left, top, width, height)
                subtitle_frame = subtitle_box.text_frame
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.text = page.subtitle
                subtitle_para.font.size = Pt(20)
                subtitle_para.font.name = fonts["body"]
                subtitle_para.font.color.rgb = self._hex_to_rgb(colors["secondary"])
                subtitle_para.alignment = PP_ALIGN.CENTER
            
            # 添加内容
            if page.content and "content" in positions:
                left, top, width, height = positions["content"]
                content_box = slide.shapes.add_textbox(left, top, width, height)
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                for i, item in enumerate(page.content):
                    if i == 0:
                        para = content_frame.paragraphs[0]
                    else:
                        para = content_frame.add_paragraph()
                    para.text = f"• {item}"
                    para.font.size = Pt(18)
                    para.font.name = fonts["body"]
                    para.font.color.rgb = self._hex_to_rgb(colors["text"])
                    para.space_after = Pt(12)
            
            # 添加讲者备注
            if page.notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = page.notes
        
        # 生成文件名
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        keywords_str = "_".join(request.keywords[:3]) if request.keywords else "presentation"
        file_name = f"{keywords_str}_{timestamp}.pptx"
        file_path = self.output_dir / file_name
        
        # 保存文件
        prs.save(str(file_path))
        
        return file_path
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """将十六进制颜色转换为RGBColor"""
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 6:
            r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
            return RGBColor(r, g, b)
        return RGBColor(0, 0, 0)
    
    async def regenerate_page(self, ppt_path: str, page_index: int, new_content: List[str]) -> bool:
        """重新生成单页"""
        try:
            if not PPTX_AVAILABLE:
                return False
            
            prs = Presentation(ppt_path)
            if page_index < len(prs.slides):
                slide = prs.slides[page_index]
                
                # 清除现有内容
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        shape.text_frame.clear()
                
                # 添加新内容
                content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(5))
                content_frame = content_box.text_frame
                
                for i, item in enumerate(new_content):
                    if i == 0:
                        para = content_frame.paragraphs[0]
                    else:
                        para = content_frame.add_paragraph()
                    para.text = f"• {item}"
                    para.font.size = Pt(18)
                
                prs.save(ppt_path)
                return True
            return False
        except Exception as e:
            print(f"Error regenerating page: {e}")
            return False
    
    async def change_style(self, ppt_path: str, new_style: PPTStyle) -> bool:
        """切换风格"""
        try:
            if not PPTX_AVAILABLE:
                return False
            
            prs = Presentation(ppt_path)
            colors = ColorScheme.get(new_style)
            
            # 这里可以实现风格切换逻辑
            # 由于python-pptx的限制，完整的风格切换比较复杂
            # 这里仅作为示例
            
            prs.save(ppt_path)
            return True
        except Exception as e:
            print(f"Error changing style: {e}")
            return False


class PPTPipeline:
    """PPT生成流水线"""
    
    def __init__(self, config: Optional[PPTConfig] = None):
        self.config = config or PPTConfig()
        self.generator = PPTGenerator()
    
    async def run(self, user_input: str, user_id: str = "default") -> PPTResult:
        """执行生成流水线"""
        request = PPTRequest(
            user_input=user_input,
            user_id=user_id,
            config=self.config
        )
        
        result = await self.generator.generate(request)
        return result
    
    async def run_with_assets(
        self,
        user_input: str,
        selected_assets: List[str],
        user_id: str = "default"
    ) -> PPTResult:
        """带素材执行生成"""
        request = PPTRequest(
            user_input=user_input,
            uploaded_images=selected_assets,
            user_id=user_id,
            config=self.config
        )
        
        return await self.generator.generate(request)
    
    def run_sync(self, user_input: str, user_id: str = "default") -> PPTResult:
        """同步执行生成"""
        return asyncio.run(self.run(user_input, user_id))


if __name__ == "__main__":
    # 示例使用
    pipeline = PPTPipeline()
    
    # 测试生成
    result = pipeline.run_sync("人工智能在医疗领域的应用与发展")
    
    if result.success:
        print(f"PPT生成成功: {result.file_path}")
        print(f"生成了 {len(result.pages)} 页")
        print(f"使用了 {len(result.assets_used)} 个素材")
    else:
        print(f"PPT生成失败: {result.error}")
